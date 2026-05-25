"""
Build the Creekside Marketing pricing slide deck as a .pptx file.

Peterson uploads the .pptx to Google Drive and Drive auto-converts it to an
editable Google Slides file.

Slide 1: "Our Pricing" -- single pricing plan card
Slide 2: "How It Scales" -- fee scaling chart + tier commentary

Simplified to single pricing plan on 2026-05-25 (Plan B Shared and Plan C Retainer removed).
Pairs with proposal_chart.py (provides the chart for slide 2) and build_docx.py.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
CHART_PATH = os.path.join(SCRIPT_DIR, "out", "proposal_chart.png")
OUT_PATH = os.path.join(SCRIPT_DIR, "out", "Pricing_Plans.pptx")
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)

# Brand colors
CREEKSIDE_BLUE = RGBColor(0x1A, 0x56, 0xDB)
DARK_TEXT = RGBColor(0x11, 0x18, 0x27)
GRAY_TEXT = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG = RGBColor(0xF9, 0xFA, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]


def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tx


def add_rect(slide, left, top, width, height, fill=WHITE, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: Our Pricing
# ═══════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(BLANK_LAYOUT)

# Title
add_text(slide1, Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.7),
         "Our Pricing",
         size=32, bold=True, color=CREEKSIDE_BLUE, align=PP_ALIGN.CENTER)

# Single centered pricing card
card_top = Inches(1.3)
card_height = Inches(5.0)
card_width = Inches(5.5)
left_margin = (prs.slide_width - card_width) / 2

# Card background
add_rect(slide1, left_margin, card_top, card_width, card_height,
         fill=LIGHT_BG, line=RGBColor(0xE5, 0xE7, 0xEB))

# Header band (blue strip at top)
header_height = Inches(0.85)
add_rect(slide1, left_margin, card_top, card_width, header_height,
         fill=CREEKSIDE_BLUE)

# Header text
add_text(slide1, left_margin, card_top + Inches(0.05),
         card_width, Inches(0.30),
         "Percentage of Ad Spend",
         size=12, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide1, left_margin, card_top + Inches(0.32),
         card_width, Inches(0.50),
         "MANAGEMENT FEE",
         size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Body content -- 4 rows
row_top = card_top + header_height + Inches(0.18)
row_height = Inches(1.0)
rows = [
    ("VARIABLE FEE", "20% up to $30k\n15% from $30k-$60k\n10% above $60k"),
    ("MINIMUM FEE", "$1,500 per platform"),
    ("MONTHLY CAP", "$15,000 maximum"),
    ("ONBOARDING", "$1,500 per platform"),
]
inner_left = left_margin + Inches(0.25)
inner_width = card_width - Inches(0.5)
for r_idx, (label, value) in enumerate(rows):
    rt = row_top + row_height * r_idx
    add_text(slide1, inner_left, rt, inner_width, Inches(0.25),
             label,
             size=9, bold=True, color=GRAY_TEXT, align=PP_ALIGN.CENTER)
    add_text(slide1, inner_left, rt + Inches(0.26),
             inner_width, Inches(0.70),
             value,
             size=12, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# Footer note about minimum
footer_top = card_top + card_height + Inches(0.2)
add_text(slide1, Inches(0.7), footer_top, Inches(11.9), Inches(0.6),
         "The $1,500 minimum fee per platform acts as a floor, covering your first "
         "$7,500 in ad spend. Once you spend over $7,500, the 20% rate naturally takes over.",
         size=11, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: How It Scales
# ═══════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(BLANK_LAYOUT)

# Title
add_text(slide2, Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.7),
         "How It Scales",
         size=32, bold=True, color=CREEKSIDE_BLUE, align=PP_ALIGN.CENTER)

# Chart (embedded from proposal_chart.py output)
if os.path.exists(CHART_PATH):
    chart_width = Inches(11.5)
    chart_left = (prs.slide_width - chart_width) / 2
    slide2.shapes.add_picture(CHART_PATH, chart_left, Inches(1.2),
                              width=chart_width)
    commentary_top = Inches(6.4)
else:
    add_text(slide2, Inches(0.5), Inches(1.2),
             Inches(12.3), Inches(0.5),
             f"[Chart missing -- run proposal_chart.py first to generate {CHART_PATH}]",
             size=12, color=GRAY_TEXT, align=PP_ALIGN.CENTER)
    commentary_top = Inches(2.0)

# Commentary
add_text(slide2, Inches(1.0), commentary_top,
         Inches(11.3), Inches(0.9),
         "Your management fee scales with your ad spend. The percentage drops "
         "at $30,000 and again at $60,000 per platform. The $15,000 monthly cap "
         "means your costs are predictable even as you grow.",
         size=13, color=DARK_TEXT, align=PP_ALIGN.CENTER)

prs.save(OUT_PATH)
print(f"Saved: {OUT_PATH}")
